"""
Document conversion service.

Runs the same class of engine the big online converters use on their servers:
  - Word / PowerPoint / Excel -> PDF   via LibreOffice headless
  - PDF -> Word (.docx)                via pdf2docx (layout-aware reconstruction)
  - PDF -> PowerPoint (.pptx)          via PyMuPDF page render + python-pptx (one slide per page)
  - PDF -> Excel (.xlsx)               via PyMuPDF table detection + openpyxl

Everything happens on the server host — files go to a per-request temp folder
and are deleted straight after the response. Website visitors install nothing.

    pip install -r requirements.txt
    python server/convert_server.py            # or: npm run server

Endpoints
    GET  /health                     -> { status, libreoffice, pdf2docx, unlock, ... }
    POST /convert/word-to-pdf        -> field "file" (.docx/.doc/.odt/.rtf/.txt) -> PDF
    POST /convert/powerpoint-to-pdf  -> field "file" (.pptx/.ppt/.odp) -> PDF
    POST /convert/excel-to-pdf       -> field "file" (.xlsx/.xls/.ods/.csv) -> PDF
    POST /convert/pdf-to-word        -> field "file" (.pdf) [+ "pages" e.g. "1-3,7"] -> .docx
    POST /convert/pdf-to-powerpoint  -> field "file" (.pdf) [+ "pages" e.g. "1-3,7"] -> .pptx
    POST /convert/pdf-to-excel       -> field "file" (.pdf) [+ "pages" e.g. "1-3,7"] -> .xlsx
    POST /pdf/unlock            -> field "file" (.pdf) [+ "password"] -> decrypted PDF
    POST /pdf/protect           -> field "file" (.pdf) + "password" (4-6 chars) -> encrypted PDF
    POST /pdf/compress          -> field "file" (.pdf) + "level" [+ "targetKb"] -> smaller PDF
    POST /pdf/fill-sign         -> JSON { pdf(base64), form_fields[], signatures[] } -> { modifiedPdf }
"""

import base64
import glob
import io
import os
import re
import shutil
import subprocess
import tempfile
import warnings

warnings.filterwarnings("ignore", message=r"The `fitz` API")

from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from werkzeug.utils import secure_filename

app = Flask(__name__)

# Dev: allow everything. Production: set ALLOWED_ORIGINS to your site's origin(s),
# comma-separated, e.g. "https://imresizer.app,https://www.imresizer.app".
_origins = os.environ.get("ALLOWED_ORIGINS", "").strip()
CORS(
    app,
    origins=[o.strip() for o in _origins.split(",") if o.strip()] or "*",
    expose_headers=[
        "X-Signatures-Removed", "Content-Disposition",
        "X-Original-Size", "X-Compressed-Size", "X-Reduction", "X-Compression-Note",
    ],
)

MAX_CONTENT_LENGTH = 100 * 1024 * 1024  # 100 MB
app.config["MAX_CONTENT_LENGTH"] = MAX_CONTENT_LENGTH
CONVERT_TIMEOUT = 180  # seconds (LibreOffice subprocess)

WORD_EXT = {".docx", ".doc", ".odt", ".rtf", ".txt", ".wpd", ".fodt"}
PPT_EXT = {".pptx", ".ppt", ".odp", ".pps", ".ppsx", ".fodp", ".key"}
EXCEL_EXT = {".xlsx", ".xls", ".ods", ".csv", ".xlsm", ".fods", ".tsv"}
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

try:
    from pdf2docx import Converter as _Pdf2DocxConverter  # noqa: N814

    HAS_PDF2DOCX = True
except Exception:  # noqa: BLE001
    _Pdf2DocxConverter = None
    HAS_PDF2DOCX = False

try:
    import pikepdf

    HAS_PIKEPDF = True
except Exception:  # noqa: BLE001
    pikepdf = None
    HAS_PIKEPDF = False

try:
    import pymupdf  # PyMuPDF
    from PIL import Image as _PILImage

    HAS_COMPRESS = True
except Exception:  # noqa: BLE001
    pymupdf = None
    _PILImage = None
    HAS_COMPRESS = False

try:
    from pptx import Presentation as _PptxPresentation
    from pptx.util import Emu as _PptxEmu

    HAS_PPTX = True
except Exception:  # noqa: BLE001
    _PptxPresentation = None
    _PptxEmu = None
    HAS_PPTX = False

try:
    from openpyxl import Workbook as _OpenpyxlWorkbook
    from openpyxl.styles import Font as _XlFont, Border as _XlBorder, Side as _XlSide, PatternFill as _XlFill
    from openpyxl.utils import get_column_letter

    _XL_SIDE = _XlSide(style="thin", color="BFBFBF")
    _XL_BOX = _XlBorder(left=_XL_SIDE, right=_XL_SIDE, top=_XL_SIDE, bottom=_XL_SIDE)
    _XL_HEADER_FILL = _XlFill("solid", fgColor="F2F2F2")
    HAS_XLSX = True
except Exception:  # noqa: BLE001
    _OpenpyxlWorkbook = _XlFont = _XlBorder = _XlSide = _XlFill = None
    get_column_letter = None
    _XL_SIDE = _XL_BOX = _XL_HEADER_FILL = None
    HAS_XLSX = False


# --------------------------------------------------------------------------- #
#  LibreOffice
# --------------------------------------------------------------------------- #
def find_soffice():
    for name in ("soffice", "soffice.com", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.com",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.com",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice",
        "/usr/local/bin/soffice",
        "/opt/libreoffice/program/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    for path in candidates:
        if os.path.isfile(path):
            return path
    return None


SOFFICE = find_soffice()


def libreoffice_convert(src_path, work_dir, target):
    """Convert src_path -> `target` (e.g. 'pdf', 'docx') inside work_dir."""
    if not SOFFICE:
        raise RuntimeError("LibreOffice is not installed on the server.")

    # A private profile per call avoids LibreOffice's single-instance lock.
    profile = os.path.join(work_dir, "lo_profile")
    os.makedirs(profile, exist_ok=True)
    profile_uri = "file:///" + profile.replace("\\", "/").lstrip("/")

    cmd = [
        SOFFICE,
        "-env:UserInstallation=" + profile_uri,
        "--headless",
        "--norestore",
        "--nolockcheck",
        "--nodefault",
        "--nofirststartwizard",
        "--convert-to",
        target,
        "--outdir",
        work_dir,
        src_path,
    ]
    proc = subprocess.run(
        cmd, capture_output=True, text=True, timeout=CONVERT_TIMEOUT, cwd=work_dir
    )
    ext = "." + target.split(":")[0]
    produced = [p for p in glob.glob(os.path.join(work_dir, "*" + ext)) if p != src_path]
    if not produced:
        raise RuntimeError(
            "LibreOffice could not convert the file. "
            + (proc.stderr or proc.stdout or "").strip()
        )
    return produced[0]


# --------------------------------------------------------------------------- #
#  helpers
# --------------------------------------------------------------------------- #
def parse_pages(spec, total=None):
    """'1-3,7' (1-indexed) -> sorted 0-indexed list. None/'' -> None (all pages)."""
    spec = (spec or "").strip()
    if not spec:
        return None
    out = set()
    for chunk in spec.split(","):
        chunk = chunk.strip()
        if not chunk:
            continue
        m = re.match(r"^(\d+)\s*-\s*(\d+)$", chunk)
        if m:
            a, b = int(m.group(1)), int(m.group(2))
            for n in range(min(a, b), max(a, b) + 1):
                out.add(n - 1)
        elif chunk.isdigit():
            out.add(int(chunk) - 1)
        else:
            raise ValueError(f"Bad page range: '{chunk}'")
    pages = sorted(n for n in out if n >= 0 and (total is None or n < total))
    return pages or None


def strip_signatures(pdf):
    """Remove digital-signature fields and their on-page widgets.

    Decrypting a PDF rewrites its bytes, which always breaks any embedded
    digital signature (the signed byte range no longer matches). If we leave
    the now-broken signature in place, viewers stamp a yellow
    "Signature Not Verified" / "?" marker over the page — which is exactly
    what happens with signed-then-encrypted docs like Aadhaar. The big online
    unlock tools drop the signature instead, leaving a clean PDF. We do the
    same. Returns the number of signatures removed.
    """
    removed = 0
    root = pdf.Root

    def objgen(obj):
        try:
            return obj.objgen
        except Exception:  # noqa: BLE001
            return None

    def is_sig_dict(obj):
        try:
            if str(obj.get("/FT")) == "/Sig":
                return True
            parent = obj.get("/Parent")
            if parent is not None and str(parent.get("/FT")) == "/Sig":
                return True
        except Exception:  # noqa: BLE001
            pass
        return False

    # 1. Signature entries in the AcroForm field tree.
    sig_ids = set()
    acro = root.get("/AcroForm")
    if acro is not None and "/Fields" in acro:
        kept_fields = []
        for f in list(acro.Fields):
            if is_sig_dict(f):
                removed += 1
                sig_ids.add(objgen(f))
                for kid in (f.get("/Kids") or []):
                    sig_ids.add(objgen(kid))
            else:
                kept_fields.append(f)
        if kept_fields:
            acro.Fields = pdf.make_indirect(pikepdf.Array(kept_fields))
            if "/SigFlags" in acro:
                del acro.SigFlags
        elif "/AcroForm" in root:
            del root.AcroForm

    # 2. Signature widget annotations on every page.
    for page in pdf.pages:
        annots = page.get("/Annots")
        if annots is None:
            continue
        original = list(annots)
        kept = [
            a for a in original
            if objgen(a) not in sig_ids and not is_sig_dict(a)
        ]
        if len(kept) != len(original):
            if kept:
                page.Annots = pdf.make_indirect(pikepdf.Array(kept))
            elif "/Annots" in page:
                del page.Annots

    return removed


# --------------------------------------------------------------------------- #
#  PDF compression (PyMuPDF) — recompress embedded images, keep text vector
# --------------------------------------------------------------------------- #
COMPRESS_LEVELS = {
    # long-edge pixel cap, JPEG quality
    "light": None,                       # lossless restructure only
    "medium": {"cap": 1700, "quality": 62},
    "strong": {"cap": 1100, "quality": 45},
    "extreme": {"cap": 850, "quality": 32},
}


def _recompress_images(doc, cap, quality):
    changed = 0
    for page in doc:
        for img in page.get_images(full=True):
            xref = img[0]
            try:
                info = doc.extract_image(xref)
                if not info or not info.get("image"):
                    continue
                pil = _PILImage.open(io.BytesIO(info["image"]))
                w, h = pil.size
                if max(w, h) > cap:
                    s = cap / float(max(w, h))
                    pil = pil.resize((max(1, int(w * s)), max(1, int(h * s))), _PILImage.LANCZOS)
                if pil.mode in ("RGBA", "P", "LA", "PA"):
                    pil = pil.convert("RGB")
                elif pil.mode not in ("RGB", "L"):
                    pil = pil.convert("RGB")
                buf = io.BytesIO()
                pil.save(buf, "JPEG", quality=quality, optimize=True, progressive=True)
                if buf.tell() < len(info["image"]):
                    page.replace_image(xref, stream=buf.getvalue())
                    changed += 1
            except Exception:  # noqa: BLE001, PERF203
                continue
    return changed


def _save_optimized(doc):
    return doc.tobytes(
        garbage=4, deflate=True, deflate_images=True, deflate_fonts=True, clean=True,
    )


def compress_pdf(raw, level="medium", target_bytes=None):
    """Return (compressed_bytes, note). Text stays selectable at every level."""
    best = raw
    note = ""

    def attempt(params):
        d = pymupdf.open(stream=raw, filetype="pdf")
        try:
            if params:
                _recompress_images(d, params["cap"], params["quality"])
            return _save_optimized(d)
        finally:
            d.close()

    order = [level] if level in COMPRESS_LEVELS else ["medium"]
    if target_bytes:
        # escalate through the stronger presets until the target is met
        seq = ["light", "medium", "strong", "extreme"]
        order = seq[seq.index(order[0]):] if order[0] in seq else seq

    for name in order:
        out = attempt(COMPRESS_LEVELS[name])
        if len(out) < len(best):
            best = out
            note = f"compressed ({name})"
        if target_bytes and len(out) <= target_bytes:
            note = f"reached target ({name})"
            return out, note

    if target_bytes and len(best) > target_bytes:
        note = "smallest without further quality loss — target not reached"
    if len(best) >= len(raw):
        return raw, "already optimised — nothing to gain"
    return best, note or "compressed"


# --------------------------------------------------------------------------- #
#  PDF -> PowerPoint — render each page and drop it onto its own slide
# --------------------------------------------------------------------------- #
PT_PER_EMU = 12700          # 1 pt = 12700 EMU (python-pptx unit)
PDF_TO_PPTX_DPI = 150       # render resolution for the page images


def pdf_to_pptx_bytes(raw, pages=None):
    """PDF bytes -> .pptx bytes. One slide per page, page rendered as a picture
    that fills the slide; the slide takes the page's own size and aspect ratio
    so nothing is stretched or cropped."""
    doc = pymupdf.open(stream=raw, filetype="pdf")
    try:
        if doc.needs_pass:
            raise RuntimeError("This PDF is password-protected. Unlock it first.")
        total = doc.page_count
        indices = pages if pages is not None else list(range(total))
        indices = [i for i in indices if 0 <= i < total]
        if not indices:
            raise RuntimeError("No pages to convert.")

        zoom = PDF_TO_PPTX_DPI / 72.0
        matrix = pymupdf.Matrix(zoom, zoom)

        prs = _PptxPresentation()
        first = doc[indices[0]].rect
        prs.slide_width = _PptxEmu(int(first.width * PT_PER_EMU))
        prs.slide_height = _PptxEmu(int(first.height * PT_PER_EMU))
        blank = prs.slide_layouts[6]

        for i in indices:
            page = doc[i]
            rect = page.rect
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            png = pix.tobytes("png")
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                io.BytesIO(png), 0, 0,
                width=_PptxEmu(int(rect.width * PT_PER_EMU)),
                height=_PptxEmu(int(rect.height * PT_PER_EMU)),
            )

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    finally:
        doc.close()


# --------------------------------------------------------------------------- #
#  PDF -> Excel — detect tables + carry over borders, bold, size, colour
# --------------------------------------------------------------------------- #
_XL_NUM = re.compile(r"^-?\$?€?£?\s?\(?-?[\d,]*\.?\d+\)?%?$")


def _xl_coerce(value):
    """Turn '1,234', '(45)', '$12.50', '7%' into real numbers; leave text alone."""
    if value is None:
        return ""
    s = str(value).strip()
    if not s or not _XL_NUM.match(s):
        return s
    negative = s.startswith("(") and s.endswith(")")
    t = s.strip("()%").lstrip("$€£").replace(",", "").strip()
    try:
        n = float(t)
        if negative:
            n = -n
        return int(n) if n.is_integer() else n
    except ValueError:
        return s


def _rect_hit(a, b):
    """Do rectangles a and b overlap (with a small tolerance)?"""
    return (a[0] < b[2] - 0.5 and a[2] > b[0] + 0.5
            and a[1] < b[3] - 0.5 and a[3] > b[1] + 0.5)


def _collect_spans(page):
    """Flat list of text spans: {text, bbox, size, bold, italic, color}."""
    out = []
    try:
        data = page.get_text("dict")
    except Exception:  # noqa: BLE001
        return out
    for block in data.get("blocks", []):
        for line in block.get("lines", []):
            for s in line.get("spans", []):
                txt = s.get("text", "")
                if not txt.strip():
                    continue
                font = (s.get("font") or "").lower()
                flags = s.get("flags", 0)
                out.append({
                    "text": txt,
                    "bbox": tuple(s.get("bbox", (0, 0, 0, 0))),
                    "size": float(s.get("size", 0) or 0),
                    "bold": bool(flags & 16) or "bold" in font or "black" in font or "heavy" in font or "semibold" in font,
                    "italic": bool(flags & 2) or "italic" in font or "oblique" in font,
                    "color": int(s.get("color", 0) or 0),
                })
    return out


def _xl_style_in(spans, bbox, median_size):
    """Aggregate the span style inside a cell bbox -> style dict or None."""
    hits = [s for s in spans if _rect_hit(s["bbox"], bbox)]
    if not hits:
        return None
    size = max(s["size"] for s in hits)
    bold = any(s["bold"] for s in hits)
    big = size >= median_size + 1.5
    colors = [s["color"] for s in hits if s["color"] not in (0,)]
    return {
        "bold": bold or big,
        "italic": any(s["italic"] for s in hits),
        "size": round(min(max(size, 6.0), 36.0), 1) if size else None,
        "color": ("%06X" % colors[0]) if colors else None,
        "big": big,
    }


def _xl_font(style):
    if not style:
        return None
    kw = {}
    if style.get("bold"):
        kw["bold"] = True
    if style.get("italic"):
        kw["italic"] = True
    if style.get("size"):
        kw["size"] = style["size"]
    if style.get("color") and style["color"].upper() not in ("000000", "FFFFFF"):
        kw["color"] = style["color"]
    return _XlFont(**kw) if kw else None


def _xl_cluster_rows(words, tol):
    words = sorted(words, key=lambda w: ((w[1] + w[3]) / 2, w[0]))
    rows, cur, cy = [], [], None
    for w in words:
        c = (w[1] + w[3]) / 2
        if cy is None or abs(c - cy) <= tol:
            cur.append(w)
            cy = c if cy is None else (cy + c) / 2
        else:
            rows.append(cur)
            cur = [w]
            cy = c
    if cur:
        rows.append(cur)
    return rows


def _xl_whitespace_table(page):
    """Last-resort: words -> grid of (text, bbox) cells using whitespace bands."""
    words = [w for w in page.get_text("words") if w[4].strip()]
    if not words:
        return None
    heights = sorted(w[3] - w[1] for w in words if w[3] > w[1])
    line_h = heights[len(heights) // 2] if heights else 10
    rows = _xl_cluster_rows(words, max(3, line_h * 0.6))

    x0 = min(w[0] for w in words)
    x1 = max(w[2] for w in words)
    span = int(x1 - x0) + 1
    occ = bytearray(span)
    for w in words:
        for i in range(max(0, int(w[0] - x0)), min(span, int(w[2] - x0) + 1)):
            occ[i] = 1

    gap_min = max(6, line_h * 1.1)
    bounds, i = [x0 - 1], 0
    while i < span:
        if occ[i] == 0:
            j = i
            while j < span and occ[j] == 0:
                j += 1
            if j - i >= gap_min:
                bounds.append(x0 + i + (j - i) / 2)
            i = j
        else:
            i += 1
    bounds.append(x1 + 1)
    ncol = len(bounds) - 1

    def col_of(xc):
        for k in range(ncol):
            if bounds[k] <= xc < bounds[k + 1]:
                return k
        return ncol - 1

    grid = []
    for r in rows:
        cells = [None] * ncol
        for w in sorted(r, key=lambda w: w[0]):
            ci = col_of((w[0] + w[2]) / 2)
            if cells[ci] is None:
                cells[ci] = [w[4], [w[0], w[1], w[2], w[3]]]
            else:
                cells[ci][0] += " " + w[4]
                bb = cells[ci][1]
                bb[0] = min(bb[0], w[0]); bb[1] = min(bb[1], w[1])
                bb[2] = max(bb[2], w[2]); bb[3] = max(bb[3], w[3])
        grid.append([(c[0], tuple(c[1])) if c else ("", None) for c in cells])

    keep = [k for k in range(ncol) if any(row[k][0].strip() for row in grid)]
    grid = [[row[k] for k in keep] for row in grid]
    grid = [row for row in grid if any(c[0].strip() for c in row)]
    if len(grid) < 2:
        return None
    xs = [c[1] for row in grid for c in row if c[1]]
    ys = [c[1] for row in grid for c in row if c[1]]
    bbox = (min(b[0] for b in xs), min(b[1] for b in ys),
            max(b[2] for b in xs), max(b[3] for b in ys))
    return {"cells": grid, "ruled": False, "bbox": bbox}


def _table_cells(table):
    """PyMuPDF Table -> list of rows, each a list of (text, bbox|None)."""
    values = table.extract()
    out = []
    for ri, row in enumerate(table.rows):
        vals = values[ri] if ri < len(values) else []
        line = []
        for ci, bb in enumerate(row.cells):
            txt = vals[ci] if ci < len(vals) and vals[ci] is not None else ""
            line.append((str(txt), tuple(bb) if bb else None))
        out.append(line)
    return out


def _xl_find_tables(page):
    """Return [{cells, ruled, bbox}] for the page (ruled first, then text, then whitespace)."""
    for kwargs, ruled in ((dict(), True), (dict(strategy="text"), False)):
        try:
            found = [t for t in page.find_tables(**kwargs).tables if len(t.rows) > 1]
        except Exception:  # noqa: BLE001
            found = []
        if found:
            return [{"cells": _table_cells(t), "ruled": ruled, "bbox": tuple(t.bbox)} for t in found]
    ws = _xl_whitespace_table(page)
    return [ws] if ws else []


def _xl_headings(spans, table_bboxes, median_size):
    """Lines of bold / oversized text that sit outside every table -> heading rows."""
    out = []
    for s in spans:
        if s["size"] < median_size + 1.5 and not s["bold"]:
            continue
        if not any(ch.isalpha() for ch in s["text"]):
            continue
        if any(_rect_hit(s["bbox"], tb) for tb in table_bboxes):
            continue
        out.append({
            "text": s["text"].strip(),
            "y": s["bbox"][1],
            "style": {
                "bold": True,
                "italic": s["italic"],
                "size": round(min(max(s["size"], 6.0), 36.0), 1) if s["size"] else None,
                "color": ("%06X" % s["color"]) if s["color"] not in (0,) else None,
                "big": s["size"] >= median_size + 1.5,
            },
        })
    # merge duplicate/adjacent fragments on the same visual line
    out.sort(key=lambda h: h["y"])
    merged = []
    for h in out:
        if merged and abs(h["y"] - merged[-1]["y"]) < 3:
            merged[-1]["text"] = (merged[-1]["text"] + " " + h["text"]).strip()
        else:
            merged.append(h)
    return merged


def pdf_to_xlsx_bytes(raw, pages=None):
    """PDF bytes -> .xlsx bytes. One worksheet per page; cell borders, bold text,
    headings, font size and colour are carried over from the page."""
    doc = pymupdf.open(stream=raw, filetype="pdf")
    try:
        if doc.needs_pass:
            raise RuntimeError("This PDF is password-protected. Unlock it first.")
        indices = pages if pages is not None else list(range(doc.page_count))
        indices = [i for i in indices if 0 <= i < doc.page_count]

        wb = _OpenpyxlWorkbook()
        wb.remove(wb.active)
        found_any = False

        for i in indices:
            page = doc[i]
            tables = _xl_find_tables(page)
            if not tables:
                continue

            spans = _collect_spans(page)
            sizes = sorted(s["size"] for s in spans if s["size"] > 0)
            median_size = sizes[len(sizes) // 2] if sizes else 11.0
            headings = _xl_headings(spans, [t["bbox"] for t in tables], median_size)

            ws = wb.create_sheet(f"Page {i + 1}")
            found_any = True
            row_no = 1
            col_w = {}

            # interleave headings and tables top-to-bottom
            items = ([("heading", h["y"], h) for h in headings]
                     + [("table", t["bbox"][1], t) for t in tables])
            items.sort(key=lambda it: it[1])

            for kind, _y, obj in items:
                if kind == "heading":
                    cell = ws.cell(row=row_no, column=1, value=obj["text"])
                    font = _xl_font(obj["style"])
                    if font:
                        cell.font = font
                    if obj["style"].get("big"):
                        ws.row_dimensions[row_no].height = max(18, obj["style"]["size"] * 1.4)
                    row_no += 2
                    continue

                rows = obj["cells"]
                ruled = obj["ruled"]
                header_is_bold = None
                for ri, row in enumerate(rows):
                    non_empty_bold = []
                    for ci, (text, bbox) in enumerate(row, start=1):
                        val = _xl_coerce(text)
                        c = ws.cell(row=row_no, column=ci)
                        if val != "":
                            c.value = val
                        if ruled:
                            c.border = _XL_BOX
                        style = _xl_style_in(spans, bbox, median_size) if bbox else None
                        font = _xl_font(style)
                        if font:
                            c.font = font
                        if val != "":
                            non_empty_bold.append(bool(style and style["bold"]))
                        if bbox:
                            w = (bbox[2] - bbox[0]) / 5.2
                            col_w[ci] = min(80, max(col_w.get(ci, 8.0), w))
                    if ri == 0:
                        header_is_bold = bool(non_empty_bold) and all(non_empty_bold)
                    if ri == 0 and header_is_bold:
                        for ci in range(1, len(row) + 1):
                            ws.cell(row=row_no, column=ci).fill = _XL_HEADER_FILL
                    row_no += 1
                row_no += 1  # spacer after the table

            for ci, w in col_w.items():
                ws.column_dimensions[get_column_letter(ci)].width = round(w, 1)

        if not found_any:
            raise RuntimeError(
                "No tables or tabular text were found in this PDF. "
                "Scanned/image-only PDFs need OCR first."
            )

        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()
    finally:
        doc.close()


def send_bytes(data, download_name, mimetype):
    return send_file(
        io.BytesIO(data),
        mimetype=mimetype,
        as_attachment=True,
        download_name=download_name,
        max_age=0,
    )


# --------------------------------------------------------------------------- #
#  routes
# --------------------------------------------------------------------------- #
@app.get("/health")
def health():
    return jsonify(
        {
            "status": "healthy" if (SOFFICE or HAS_PDF2DOCX or HAS_PIKEPDF) else "degraded",
            "libreoffice": bool(SOFFICE),
            "soffice": SOFFICE,
            "pdf2docx": HAS_PDF2DOCX,
            "pikepdf": HAS_PIKEPDF,
            "python_pptx": HAS_PPTX,
            "openpyxl": HAS_XLSX,
            "word_to_pdf": bool(SOFFICE),
            "powerpoint_to_pdf": bool(SOFFICE),
            "excel_to_pdf": bool(SOFFICE),
            "pdf_to_word": HAS_PDF2DOCX,
            "pdf_to_powerpoint": HAS_PPTX and HAS_COMPRESS,
            "pdf_to_excel": HAS_XLSX and HAS_COMPRESS,
            "unlock": HAS_PIKEPDF,
            "protect": HAS_PIKEPDF,
            "compress": HAS_COMPRESS,
            "fill_sign": HAS_COMPRESS,
        }
    )


@app.post("/convert/word-to-pdf")
def word_to_pdf():
    if "file" not in request.files:
        return jsonify({"error": "No file provided (expected multipart field 'file')."}), 400

    upload = request.files["file"]
    name = secure_filename(upload.filename or "document")
    ext = os.path.splitext(name)[1].lower()
    if ext not in WORD_EXT:
        return jsonify({"error": f"Unsupported file type '{ext or '?'}'."}), 415
    if not SOFFICE:
        return jsonify({"error": "LibreOffice is not installed on the conversion server."}), 503

    work_dir = tempfile.mkdtemp(prefix="imr_w2p_")
    try:
        src_path = os.path.join(work_dir, "input" + ext)
        upload.save(src_path)
        out_path = libreoffice_convert(src_path, work_dir, "pdf:writer_pdf_Export")
        with open(out_path, "rb") as fh:
            data = fh.read()
    except subprocess.TimeoutExpired:
        return jsonify({"error": "Conversion timed out. The document may be too large."}), 504
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": str(exc)}), 500
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)

    return send_bytes(data, os.path.splitext(name)[0] + ".pdf", "application/pdf")


@app.post("/convert/powerpoint-to-pdf")
def powerpoint_to_pdf():
    if "file" not in request.files:
        return jsonify({"error": "No file provided (expected multipart field 'file')."}), 400

    upload = request.files["file"]
    name = secure_filename(upload.filename or "presentation")
    ext = os.path.splitext(name)[1].lower()
    if ext not in PPT_EXT:
        return jsonify({"error": f"Unsupported file type '{ext or '?'}'. Upload a .pptx, .ppt or .odp."}), 415
    if not SOFFICE:
        return jsonify({"error": "LibreOffice is not installed on the conversion server."}), 503

    work_dir = tempfile.mkdtemp(prefix="imr_p2pdf_")
    try:
        src_path = os.path.join(work_dir, "input" + ext)
        upload.save(src_path)
        out_path = libreoffice_convert(src_path, work_dir, "pdf:impress_pdf_Export")
        with open(out_path, "rb") as fh:
            data = fh.read()
    except subprocess.TimeoutExpired:
        return jsonify({"error": "Conversion timed out. The presentation may be too large."}), 504
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": str(exc)}), 500
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)

    return send_bytes(data, os.path.splitext(name)[0] + ".pdf", "application/pdf")


@app.post("/convert/excel-to-pdf")
def excel_to_pdf():
    if "file" not in request.files:
        return jsonify({"error": "No file provided (expected multipart field 'file')."}), 400

    upload = request.files["file"]
    name = secure_filename(upload.filename or "spreadsheet")
    ext = os.path.splitext(name)[1].lower()
    if ext not in EXCEL_EXT:
        return jsonify({"error": f"Unsupported file type '{ext or '?'}'. Upload a .xlsx, .xls, .ods or .csv."}), 415
    if not SOFFICE:
        return jsonify({"error": "LibreOffice is not installed on the conversion server."}), 503

    work_dir = tempfile.mkdtemp(prefix="imr_x2pdf_")
    try:
        src_path = os.path.join(work_dir, "input" + ext)
        upload.save(src_path)
        out_path = libreoffice_convert(src_path, work_dir, "pdf:calc_pdf_Export")
        with open(out_path, "rb") as fh:
            data = fh.read()
    except subprocess.TimeoutExpired:
        return jsonify({"error": "Conversion timed out. The spreadsheet may be too large."}), 504
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": str(exc)}), 500
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)

    return send_bytes(data, os.path.splitext(name)[0] + ".pdf", "application/pdf")


@app.post("/convert/pdf-to-word")
def pdf_to_word():
    if "file" not in request.files:
        return jsonify({"error": "No file provided (expected multipart field 'file')."}), 400

    upload = request.files["file"]
    name = secure_filename(upload.filename or "document.pdf")
    if os.path.splitext(name)[1].lower() != ".pdf":
        return jsonify({"error": "Please upload a .pdf file."}), 415
    if not HAS_PDF2DOCX:
        return jsonify({"error": "The PDF-to-Word engine (pdf2docx) is not installed on the server."}), 503

    try:
        pages = parse_pages(request.form.get("pages"))
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400

    work_dir = tempfile.mkdtemp(prefix="imr_p2w_")
    try:
        pdf_path = os.path.join(work_dir, "input.pdf")
        docx_path = os.path.join(work_dir, "output.docx")
        upload.save(pdf_path)

        cv = _Pdf2DocxConverter(pdf_path)
        try:
            cv.convert(docx_path, pages=pages)
        finally:
            cv.close()

        if not os.path.isfile(docx_path) or os.path.getsize(docx_path) == 0:
            raise RuntimeError("The converter produced an empty document.")
        with open(docx_path, "rb") as fh:
            data = fh.read()
    except Exception as exc:  # noqa: BLE001
        msg = str(exc)
        if "password" in msg.lower() or "encrypt" in msg.lower():
            msg = "This PDF is password-protected. Unlock it first."
        elif not msg:
            msg = "Could not convert this PDF."
        return jsonify({"error": msg}), 500
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)

    return send_bytes(data, os.path.splitext(name)[0] + ".docx", DOCX_MIME)


@app.post("/convert/pdf-to-powerpoint")
def pdf_to_powerpoint():
    if "file" not in request.files:
        return jsonify({"error": "No file provided (expected multipart field 'file')."}), 400

    upload = request.files["file"]
    name = secure_filename(upload.filename or "document.pdf")
    if os.path.splitext(name)[1].lower() != ".pdf":
        return jsonify({"error": "Please upload a .pdf file."}), 415
    if not (HAS_PPTX and HAS_COMPRESS):
        return jsonify({"error": "The PDF-to-PowerPoint engine (python-pptx / PyMuPDF) is not installed on the server."}), 503

    try:
        pages = parse_pages(request.form.get("pages"))
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400

    raw = upload.read()
    try:
        data = pdf_to_pptx_bytes(raw, pages=pages)
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": str(exc) or "Could not convert this PDF."}), 500

    return send_bytes(data, os.path.splitext(name)[0] + ".pptx", PPTX_MIME)


@app.post("/convert/pdf-to-excel")
def pdf_to_excel():
    if "file" not in request.files:
        return jsonify({"error": "No file provided (expected multipart field 'file')."}), 400

    upload = request.files["file"]
    name = secure_filename(upload.filename or "document.pdf")
    if os.path.splitext(name)[1].lower() != ".pdf":
        return jsonify({"error": "Please upload a .pdf file."}), 415
    if not (HAS_XLSX and HAS_COMPRESS):
        return jsonify({"error": "The PDF-to-Excel engine (openpyxl / PyMuPDF) is not installed on the server."}), 503

    try:
        pages = parse_pages(request.form.get("pages"))
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400

    raw = upload.read()
    try:
        data = pdf_to_xlsx_bytes(raw, pages=pages)
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": str(exc) or "Could not convert this PDF."}), 500

    return send_bytes(data, os.path.splitext(name)[0] + ".xlsx", XLSX_MIME)


@app.post("/pdf/unlock")
def pdf_unlock():
    if "file" not in request.files:
        return jsonify({"error": "No file provided (expected multipart field 'file')."}), 400

    upload = request.files["file"]
    name = secure_filename(upload.filename or "document.pdf")
    if os.path.splitext(name)[1].lower() != ".pdf":
        return jsonify({"error": "Please upload a .pdf file."}), 415
    if not HAS_PIKEPDF:
        return jsonify({"error": "The unlock engine (pikepdf) is not installed on the server."}), 503

    password = request.form.get("password", "") or ""
    raw = upload.read()

    try:
        # An owner-password-only PDF (restricted printing/copying but opens
        # freely) unlocks with an empty password. A PDF that needs a password
        # just to open needs the real one.
        src = pikepdf.open(io.BytesIO(raw), password=password)
    except pikepdf.PasswordError:
        if not password:
            return jsonify({"error": "This PDF is protected — enter its password.", "needsPassword": True}), 401
        return jsonify({"error": "Wrong password. Check it and try again.", "needsPassword": True}), 401
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": f"Could not read this PDF: {exc}"}), 400

    # A signed PDF's signature can't survive decryption (the byte range breaks),
    # so viewers would stamp "Signature Not Verified" over it. Strip it for a
    # clean result — unless the caller explicitly opts to keep it.
    keep_sig = (request.form.get("keepSignature", "") or "").lower() in ("1", "true", "yes")
    sigs_removed = 0
    if not keep_sig:
        try:
            sigs_removed = strip_signatures(src)
        except Exception:  # noqa: BLE001
            sigs_removed = 0

    try:
        out = io.BytesIO()
        # save() with no encryption argument writes a fully decrypted PDF with
        # every restriction removed.
        src.save(out, fix_metadata_version=True)
        src.close()
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": f"Could not remove the protection: {exc}"}), 500

    resp = send_bytes(out.getvalue(), os.path.splitext(name)[0] + "-unlocked.pdf", "application/pdf")
    resp.headers["X-Signatures-Removed"] = str(sigs_removed)
    return resp


@app.post("/pdf/protect")
def pdf_protect():
    if "file" not in request.files:
        return jsonify({"error": "No file provided (expected multipart field 'file')."}), 400

    upload = request.files["file"]
    name = secure_filename(upload.filename or "document.pdf")
    if os.path.splitext(name)[1].lower() != ".pdf":
        return jsonify({"error": "Please upload a .pdf file."}), 415
    if not HAS_PIKEPDF:
        return jsonify({"error": "The protect engine (pikepdf) is not installed on the server."}), 503

    password = request.form.get("password", "") or ""
    if not (4 <= len(password) <= 6):
        return jsonify({"error": "The password must be 4 to 6 characters."}), 400

    raw = upload.read()
    try:
        src = pikepdf.open(io.BytesIO(raw))
    except pikepdf.PasswordError:
        return jsonify({"error": "This PDF is already password-protected. Unlock it first."}), 400
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": f"Could not read this PDF: {exc}"}), 400

    try:
        out = io.BytesIO()
        # Same password to open and to change permissions. AES-256 (R=6).
        src.save(
            out,
            encryption=pikepdf.Encryption(owner=password, user=password, R=6),
            fix_metadata_version=True,
        )
        src.close()
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": f"Could not protect this PDF: {exc}"}), 500

    return send_bytes(out.getvalue(), os.path.splitext(name)[0] + "-protected.pdf", "application/pdf")


@app.post("/pdf/compress")
def pdf_compress():
    if "file" not in request.files:
        return jsonify({"error": "No file provided (expected multipart field 'file')."}), 400

    upload = request.files["file"]
    name = secure_filename(upload.filename or "document.pdf")
    if os.path.splitext(name)[1].lower() != ".pdf":
        return jsonify({"error": "Please upload a .pdf file."}), 415
    if not HAS_COMPRESS:
        return jsonify({"error": "The compression engine (PyMuPDF) is not installed on the server."}), 503

    level = (request.form.get("level", "medium") or "medium").lower()
    if level not in COMPRESS_LEVELS:
        level = "medium"

    target_bytes = None
    spec = (request.form.get("targetKb", "") or "").strip()
    if spec:
        try:
            target_bytes = max(1, int(float(spec) * 1024))
        except ValueError:
            return jsonify({"error": f"Invalid target size: {spec}"}), 400

    raw = upload.read()
    original = len(raw)
    try:
        out, note = compress_pdf(raw, level=level, target_bytes=target_bytes)
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": f"Could not compress this PDF: {exc}"}), 500

    reduction = round((1 - len(out) / original) * 100, 1) if original and len(out) < original else 0.0
    resp = send_bytes(out, os.path.splitext(name)[0] + "-compressed.pdf", "application/pdf")
    resp.headers["X-Original-Size"] = str(original)
    resp.headers["X-Compressed-Size"] = str(len(out))
    resp.headers["X-Reduction"] = str(reduction)
    resp.headers["X-Compression-Note"] = note
    return resp


@app.post("/pdf/fill-sign")
def pdf_fill_sign():
    """Stamp typed text and signature images onto a PDF (PyMuPDF).

    JSON body:
        pdf          base64 of the source PDF
        form_fields  [{ page, x, y, value, fontsize? }]        (page is 0-based)
        signatures   [{ page, x, y, width, height, image }]    (image = data URL)
    Returns { status, modifiedPdf(base64) }.
    """
    if not HAS_COMPRESS:
        return jsonify({"error": "The PDF engine (PyMuPDF) is not installed on the server."}), 503

    data = request.get_json(silent=True) or {}
    if "pdf" not in data:
        return jsonify({"error": "No PDF data provided."}), 400
    try:
        pdf_bytes = base64.b64decode(data["pdf"])
    except Exception:  # noqa: BLE001
        return jsonify({"error": "Invalid PDF data."}), 400

    try:
        doc = pymupdf.open(stream=pdf_bytes, filetype="pdf")
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": f"Could not read this PDF: {exc}"}), 400

    try:
        for field in data.get("form_fields", []) or []:
            page_no = int(field.get("page", 0))
            if 0 <= page_no < doc.page_count:
                doc[page_no].insert_text(
                    (float(field["x"]), float(field["y"])),
                    str(field.get("value", "")),
                    fontsize=float(field.get("fontsize", 12)),
                    color=(0, 0, 0),
                )

        for sig in data.get("signatures", []) or []:
            page_no = int(sig.get("page", 0))
            if not (0 <= page_no < doc.page_count):
                continue
            src = sig.get("image", "")
            raw = base64.b64decode(src.split(",", 1)[1] if "," in src else src)
            img = _PILImage.open(io.BytesIO(raw))
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            x, y = float(sig["x"]), float(sig["y"])
            doc[page_no].insert_image(
                pymupdf.Rect(x, y, x + float(sig["width"]), y + float(sig["height"])),
                stream=buf.getvalue(),
            )

        out = doc.tobytes(deflate=True)
        doc.close()
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": f"Could not process this PDF: {exc}"}), 500

    return jsonify({
        "status": "success",
        "message": "PDF filled and signed successfully",
        "modifiedPdf": base64.b64encode(out).decode("ascii"),
    })


if __name__ == "__main__":
    print(f"LibreOffice : {SOFFICE or 'NOT FOUND'}")
    print(f"pdf2docx    : {'ok' if HAS_PDF2DOCX else 'NOT INSTALLED'}")
    print(f"pikepdf     : {'ok' if HAS_PIKEPDF else 'NOT INSTALLED'}")
    print(f"python-pptx : {'ok' if HAS_PPTX else 'NOT INSTALLED'}")
    print(f"openpyxl    : {'ok' if HAS_XLSX else 'NOT INSTALLED'}")
    # localhost only, no debugger — this process shells out to converters.
    app.run(host="127.0.0.1", port=int(os.environ.get("PORT", 5000)), debug=False)
